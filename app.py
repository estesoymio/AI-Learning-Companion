import hashlib
import html
import json
import os
import re
import time
from io import BytesIO
from pathlib import Path

import fitz
import streamlit as st
from dotenv import load_dotenv
from docx import Document as DocxDocument
from groq import APIConnectionError, AuthenticationError, RateLimitError
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_core.documents import Document
from langchain_core.prompts import ChatPromptTemplate
from langchain_groq import ChatGroq
from langchain_openai import OpenAIEmbeddings
from pptx import Presentation


load_dotenv()

openai_api_key = os.getenv("OPENAI_API_KEY")
groq_api_key = os.getenv("GROQ_API_KEY")

if openai_api_key:
    os.environ["OPENAI_API_KEY"] = openai_api_key

if groq_api_key:
    os.environ["GROQ_API_KEY"] = groq_api_key

llm = (
    ChatGroq(groq_api_key=groq_api_key, model_name="llama-3.1-8b-instant")
    if groq_api_key
    else None
)

prompt = ChatPromptTemplate.from_template(
    """
You are an expert AI learning companion helping students understand textbook material.

Your primary goal is to TEACH, not summarize.

Rules:
1. Answer using only the information in the provided context.
2. Use the conversation history to resolve follow-up questions.
3. If the answer is not supported by the context, clearly state that you could not find enough information in the provided material.
4. NEVER mention figure numbers, table numbers, page numbers, or diagram references (e.g., "Figure 14-22", "Table 3.1", "as shown above").
5. If the context references a figure or diagram that is unavailable, rewrite the explanation so it is completely self-contained and describe the concept in words instead.
6. Do NOT say phrases like "we don't have the figure", "the image is missing", or "as shown in the diagram".
7. When explaining technical concepts:
   - Start with a simple definition.
   - Explain the idea step by step.
   - Include an example if the user asks for one or if it improves understanding.
   - Use bullet points or numbered lists when appropriate.
8. Write naturally, as if teaching from memory rather than quoting a textbook.

Conversation history:
{chat_history}

<context>
{context}
</context>

Latest question:
{input}
"""
)

quiz_prompt = ChatPromptTemplate.from_template(
    """
    Create a {difficulty} multiple-choice quiz about the topic "{topic}".
    Use only the provided context. Do not add facts that are not supported
    by the context. Create exactly {question_count} questions.

    Return only a valid JSON array. Do not use Markdown or code fences.
    Each array item must have this exact structure:
    {{
      "question": "Question text",
      "options": {{
        "A": "First option",
        "B": "Second option",
        "C": "Third option",
        "D": "Fourth option"
      }},
      "correct_answer": "A",
      "explanation": "Why the correct answer is supported by the context"
    }}

    The correct_answer must be exactly A, B, C, or D. Make every question
    answerable from the context and give one unambiguously correct option.
    If there is not enough relevant context for the topic, say so instead of
    creating unsupported questions.

    <context>
    {context}
    </context>
    """
)

flashcard_prompt = ChatPromptTemplate.from_template(
    """
    Create exactly {card_count} {detail_level} flashcards about the topic "{topic}".
    Use only the provided context. Each flashcard should help a student study an
    important idea, definition, process, example, or relationship from the topic.
    Do not add facts that are not supported by the context.

    Format each card exactly like this:
    CARD 1
    FRONT: Question or term
    BACK: Clear answer or explanation

    CARD 2
    FRONT: Question or term
    BACK: Clear answer or explanation

    If there is not enough relevant context for the topic, say so instead of
    creating unsupported flashcards.

    <context>
    {context}
    </context>
    """
)


SUPPORTED_FILE_TYPES = ["pdf", "docx", "pptx", "txt", "md", "csv"]
CONTENTS_HEADING_PATTERN = re.compile(
    r"^\s*(table\s+of\s+contents|contents)\s*$",
    re.IGNORECASE,
)
CONTENTS_ENTRY_PATTERN = re.compile(
    r"^\s*(?:\d+(?:\.\d+)*[.)]?\s+)?(.+?)"
    r"(?:\s*\.{2,}\s*|\s{2,})(\d{1,4})\s*$"
)


def uploaded_file_signature(uploaded_files):
    return tuple(
        (
            uploaded_file.name,
            hashlib.sha256(uploaded_file.getvalue()).hexdigest(),
        )
        for uploaded_file in uploaded_files
    )


def clean_topic_title(title):
    return re.sub(r"\s+", " ", title).strip(" .\t-")


def unique_topics(topics):
    unique = []
    seen = set()

    for topic in topics:
        title = clean_topic_title(topic["title"])
        normalized = title.casefold()
        if not title or normalized in seen:
            continue
        seen.add(normalized)
        unique.append({**topic, "title": title})

    return unique


def topics_from_contents_text(text):
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    heading_indexes = [
        index
        for index, line in enumerate(lines)
        if CONTENTS_HEADING_PATTERN.match(line)
    ]
    topics = []

    for heading_index in heading_indexes:
        for line in lines[heading_index + 1:heading_index + 80]:
            match = CONTENTS_ENTRY_PATTERN.match(line)
            if match:
                topics.append(
                    {
                        "title": match.group(1),
                        "level": 1,
                        "page": int(match.group(2)),
                    }
                )
            elif topics and len(line) > 100:
                break

    return unique_topics(topics)


def extract_file_topics(uploaded_file):
    file_bytes = uploaded_file.getvalue()
    extension = Path(uploaded_file.name).suffix.lower()

    if extension == ".pdf":
        with fitz.open(stream=file_bytes, filetype="pdf") as pdf:
            bookmarks = [
                {"level": level, "title": title, "page": page}
                for level, title, page, *_ in pdf.get_toc(simple=False)
            ]
            if bookmarks:
                return unique_topics(bookmarks)

            opening_text = "\n".join(
                page.get_text() for page in list(pdf)[:15]
            )
            return topics_from_contents_text(opening_text)

    if extension == ".docx":
        document = DocxDocument(BytesIO(file_bytes))
        toc_topics = []
        heading_topics = []

        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            style_name = paragraph.style.name.lower() if paragraph.style else ""
            if not text:
                continue
            if style_name.startswith("toc"):
                match = re.match(r"toc\s*(\d+)", style_name)
                toc_topics.append(
                    {
                        "title": re.sub(r"\s+\d+\s*$", "", text),
                        "level": int(match.group(1)) if match else 1,
                    }
                )
            elif style_name.startswith("heading"):
                match = re.match(r"heading\s*(\d+)", style_name)
                heading_topics.append(
                    {
                        "title": text,
                        "level": int(match.group(1)) if match else 1,
                    }
                )

        return unique_topics(toc_topics or heading_topics)

    if extension == ".pptx":
        presentation = Presentation(BytesIO(file_bytes))
        topics = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            title_shape = slide.shapes.title
            if title_shape and title_shape.text.strip():
                topics.append(
                    {
                        "title": title_shape.text,
                        "level": 1,
                        "slide": slide_number,
                    }
                )
        return unique_topics(topics)

    text = file_bytes.decode("utf-8-sig", errors="replace")
    if extension == ".md":
        markdown_topics = []
        for line in text.splitlines():
            match = re.match(r"^(#{1,6})\s+(.+?)\s*#*\s*$", line)
            if match:
                markdown_topics.append(
                    {"title": match.group(2), "level": len(match.group(1))}
                )
        return unique_topics(markdown_topics)

    return topics_from_contents_text(text)


def extract_uploaded_topics(uploaded_files):
    topics_by_file = {}
    errors = []

    for uploaded_file in uploaded_files:
        try:
            topics_by_file[uploaded_file.name] = extract_file_topics(uploaded_file)
        except Exception as error:
            topics_by_file[uploaded_file.name] = []
            errors.append((uploaded_file.name, str(error)))

    return topics_by_file, errors


def topic_label(topic):
    location = ""
    if "page" in topic:
        location = f" (page {topic['page']})"
    elif "slide" in topic:
        location = f" (slide {topic['slide']})"
    return f"{'  ' * max(topic.get('level', 1) - 1, 0)}- {topic['title']}{location}"


def text_document(text, source, **metadata):
    return Document(
        page_content=text,
        metadata={"source": source, "chapter": Path(source).stem, **metadata},
    )


def load_uploaded_file(uploaded_file):
    file_bytes = uploaded_file.getvalue()
    extension = Path(uploaded_file.name).suffix.lower()

    if extension == ".pdf":
        documents = []
        with fitz.open(stream=file_bytes, filetype="pdf") as pdf:
            for page in pdf:
                text = page.get_text()
                if text.strip():
                    documents.append(
                        text_document(
                            text,
                            uploaded_file.name,
                            page=page.number + 1,
                        )
                    )
        return documents

    if extension == ".docx":
        document = DocxDocument(BytesIO(file_bytes))
        text = "\n".join(paragraph.text for paragraph in document.paragraphs)
        return [text_document(text, uploaded_file.name)] if text.strip() else []

    if extension == ".pptx":
        presentation = Presentation(BytesIO(file_bytes))
        documents = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            text = "\n".join(
                shape.text for shape in slide.shapes if hasattr(shape, "text")
            )
            if text.strip():
                documents.append(
                    text_document(text, uploaded_file.name, slide=slide_number)
                )
        return documents

    text = file_bytes.decode("utf-8-sig", errors="replace")
    return [text_document(text, uploaded_file.name)] if text.strip() else []


def load_uploaded_documents(uploaded_files):
    documents = []
    failed_files = []

    for uploaded_file in uploaded_files:
        try:
            documents.extend(load_uploaded_file(uploaded_file))
        except Exception as error:
            failed_files.append((uploaded_file.name, str(error)))

    return documents, failed_files


def create_vector_embedding(uploaded_files):
    docs, failed_files = load_uploaded_documents(uploaded_files)

    for filename, error in failed_files:
        st.error(f"Could not read {filename}: {error}")

    if not docs:
        st.error("No readable text was found in the uploaded files.")
        return

    selected_source = st.session_state.selected_source
    if selected_source != "All Files":
        docs = [
            doc for doc in docs
            if doc.metadata.get("source") == selected_source
        ]

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200
    )
    final_documents = text_splitter.split_documents(docs)

    st.session_state.embeddings = OpenAIEmbeddings()
    st.session_state.final_documents = final_documents
    st.session_state.vectors = FAISS.from_documents(
        final_documents,
        st.session_state.embeddings
    )
    st.session_state.indexed_source = selected_source
    st.session_state.indexed_files = uploaded_file_signature(uploaded_files)
    st.session_state.chat_messages = []
    st.session_state.pop("quiz_response", None)
    st.session_state.pop("flashcard_response", None)
    for key in (
        "quiz_question_index",
        "quiz_score",
        "quiz_answered",
        "quiz_selected_answer",
    ):
        st.session_state.pop(key, None)


def vector_database_ready():
    if "vectors" not in st.session_state:
        st.warning("Upload study files and click 'Prepare Learning Material' first.")
        return False

    if st.session_state.get("indexed_source") != st.session_state.selected_source:
        st.warning("The selected file changed. Prepare the learning material again.")
        return False

    return True


def parse_quiz_response(quiz_response):
    cleaned_response = quiz_response.strip()
    if cleaned_response.startswith("```"):
        cleaned_response = re.sub(r"^```(?:json)?\s*", "", cleaned_response)
        cleaned_response = re.sub(r"\s*```$", "", cleaned_response)

    try:
        parsed = json.loads(cleaned_response)
    except json.JSONDecodeError:
        return []

    if not isinstance(parsed, list):
        return []

    questions = []
    for item in parsed:
        if not isinstance(item, dict):
            continue
        options = item.get("options")
        correct_answer = str(item.get("correct_answer", "")).upper()
        if (
            not item.get("question")
            or not isinstance(options, dict)
            or any(option not in options for option in "ABCD")
            or correct_answer not in "ABCD"
        ):
            continue
        questions.append(
            {
                "question": str(item["question"]),
                "options": {key: str(options[key]) for key in "ABCD"},
                "correct_answer": correct_answer,
                "explanation": str(item.get("explanation", "")),
            }
        )

    return questions


def reset_quiz_state():
    st.session_state.quiz_question_index = 0
    st.session_state.quiz_score = 0
    st.session_state.quiz_answered = False
    st.session_state.quiz_selected_answer = None


def render_interactive_quiz(quiz_response):
    questions = quiz_response["questions"]
    question_index = st.session_state.get("quiz_question_index", 0)

    if question_index >= len(questions):
        st.success(
            f"Quiz complete. Your score is {st.session_state.quiz_score} "
            f"out of {len(questions)}."
        )
        st.progress(1.0)
        if st.button("Restart Quiz", use_container_width=True):
            reset_quiz_state()
            st.rerun()
        return

    question = questions[question_index]
    st.caption(f"Question {question_index + 1} of {len(questions)}")
    st.progress(question_index / len(questions))
    st.markdown(f"### {question['question']}")

    option_labels = [
        f"{key}. {question['options'][key]}"
        for key in "ABCD"
    ]
    selected_label = st.radio(
        "Choose one answer",
        option_labels,
        index=None,
        key=f"quiz_choice_{quiz_response['quiz_id']}_{question_index}",
        disabled=st.session_state.get("quiz_answered", False),
    )

    if not st.session_state.get("quiz_answered", False):
        if st.button(
            "Check Answer",
            type="primary",
            use_container_width=True,
            disabled=selected_label is None,
        ):
            selected_answer = selected_label[0]
            st.session_state.quiz_selected_answer = selected_answer
            st.session_state.quiz_answered = True
            if selected_answer == question["correct_answer"]:
                st.session_state.quiz_score += 1
            st.rerun()
        return

    selected_answer = st.session_state.quiz_selected_answer
    correct_answer = question["correct_answer"]
    correct_text = question["options"][correct_answer]
    if selected_answer == correct_answer:
        st.success("Correct!")
    else:
        st.error(
            f"Incorrect. The correct answer is {correct_answer}. {correct_text}"
        )

    if question["explanation"]:
        st.info(question["explanation"])

    if st.button("Next Question", type="primary", use_container_width=True):
        st.session_state.quiz_question_index += 1
        st.session_state.quiz_answered = False
        st.session_state.quiz_selected_answer = None
        st.rerun()


def parse_flashcards(flashcard_response):
    cards = []

    for block in flashcard_response.split("CARD ")[1:]:
        if "FRONT:" not in block or "BACK:" not in block:
            continue

        front_and_back = block.split("FRONT:", maxsplit=1)[1]
        front, back = front_and_back.split("BACK:", maxsplit=1)
        cards.append({"front": front.strip(), "back": back.strip()})

    return cards


def render_flashcard(card, index, deck_id):
    state_key = f"flashcard_flipped_{deck_id}_{index}"
    if state_key not in st.session_state:
        st.session_state[state_key] = False

    flipped = st.session_state[state_key]
    side = "Answer" if flipped else "Question"
    content = card["back"] if flipped else card["front"]
    safe_content = html.escape(content).replace("\n", "<br>")
    face_class = "flashcard-back" if flipped else "flashcard-front"

    st.markdown(
        f"""
        <div class="flashcard {face_class}">
            <div class="flashcard-number">CARD {index}</div>
            <div class="flashcard-side">{side}</div>
            <div class="flashcard-content">{safe_content}</div>
            <div class="flashcard-hint">Flip the card to see the other side</div>
        </div>
        """,
        unsafe_allow_html=True,
    )

    if st.button(
        "Show question" if flipped else "Show answer",
        key=f"flip_{deck_id}_{index}",
        use_container_width=True,
    ):
        st.session_state[state_key] = not flipped
        st.rerun()


def format_chat_history(messages):
    recent_messages = messages[-8:]
    return "\n".join(
        f"{message['role'].title()}: {message['content']}"
        for message in recent_messages
    )


def source_caption(doc):
    source = doc.metadata.get("source", "Unknown file")
    if "page" in doc.metadata:
        return f"Source: {source}, page {doc.metadata['page']}"
    if "slide" in doc.metadata:
        return f"Source: {source}, slide {doc.metadata['slide']}"
    return f"Source: {source}"


def ai_provider_ready():
    if llm is None:
        st.error(
            "GROQ_API_KEY is missing. Add a valid Groq API key to .env "
            "and restart the app."
        )
        return False

    return True


def invoke_ai_chain(chain, inputs):

    try:
        return chain.invoke(inputs)
    except AuthenticationError:
        st.error(
            "Groq rejected the API key. Generate a new key, update "
            "GROQ_API_KEY in .env, and restart Streamlit."
        )
    except RateLimitError:
        st.error("The Groq rate limit was reached. Please wait and try again.")
    except APIConnectionError:
        st.error("Could not connect to Groq. Check your connection and try again.")
    except Exception as error:
        st.error(f"The AI request failed: {error}")

    return None


def answer_chat_question(user_prompt):
    if not ai_provider_ready():
        return None

    chat_history = format_chat_history(st.session_state.chat_messages[:-1])
    retrieval_query = f"{chat_history}\nStudent: {user_prompt}".strip()
    retriever = st.session_state.vectors.as_retriever(search_kwargs={"k": 4})
    context_docs = retriever.invoke(retrieval_query)
    document_chain = create_stuff_documents_chain(llm, prompt)

    start = time.process_time()
    answer = invoke_ai_chain(
        document_chain,
        {
            "input": user_prompt,
            "chat_history": chat_history or "No earlier messages.",
            "context": context_docs,
        },
    )

    if answer is None:
        return None

    return {
        "content": answer,
        "context": context_docs,
        "response_time": time.process_time() - start,
    }


st.title("AI Learning Companion")
st.caption("Upload your own material, then chat, practise with quizzes, or revise with flashcards.")

st.markdown(
    """
    <style>
    .flashcard {
        min-height: 280px;
        border-radius: 18px;
        padding: 1.5rem;
        margin: 0.5rem 0 0.75rem;
        display: flex;
        flex-direction: column;
        justify-content: center;
        text-align: center;
        box-shadow: 0 8px 24px rgba(15, 23, 42, 0.14);
        border: 1px solid rgba(100, 116, 139, 0.25);
    }
    .flashcard-front {
        background: linear-gradient(145deg, #eff6ff, #dbeafe);
        color: #172554;
    }
    .flashcard-back {
        background: linear-gradient(145deg, #ecfdf5, #d1fae5);
        color: #052e16;
    }
    .flashcard-number {
        font-size: 0.75rem;
        font-weight: 700;
        letter-spacing: 0.12em;
        opacity: 0.65;
    }
    .flashcard-side {
        margin-top: 0.5rem;
        font-size: 0.9rem;
        font-weight: 700;
        text-transform: uppercase;
    }
    .flashcard-content {
        margin: 1.25rem 0;
        font-size: 1.2rem;
        line-height: 1.6;
        font-weight: 600;
    }
    .flashcard-hint {
        font-size: 0.78rem;
        opacity: 0.65;
    }
    </style>
    """,
    unsafe_allow_html=True,
)

if "chat_messages" not in st.session_state:
    st.session_state.chat_messages = []

with st.sidebar:
    st.header("Study Material")
    uploaded_files = st.file_uploader(
        "Upload one or more files",
        type=SUPPORTED_FILE_TYPES,
        accept_multiple_files=True,
        help="Supported formats: PDF, Word, PowerPoint, TXT, Markdown, and CSV.",
    )

    current_files = uploaded_file_signature(uploaded_files)
    if st.session_state.get("indexed_files") not in (None, current_files):
        for key in ("vectors", "final_documents", "indexed_source", "indexed_files"):
            st.session_state.pop(key, None)

    source_options = ["All Files"] + [
        uploaded_file.name for uploaded_file in uploaded_files
    ]
    if st.session_state.get("selected_source") not in source_options:
        st.session_state.selected_source = "All Files"

    st.selectbox(
        "Material to study",
        source_options,
        key="selected_source"
    )

    if st.button(
        "Prepare Learning Material",
        use_container_width=True,
        disabled=not uploaded_files,
    ):
        create_vector_embedding(uploaded_files)

        if "vectors" in st.session_state:
            st.success(f"Ready for {st.session_state.indexed_source}.")
        else:
            st.error("The learning material could not be prepared.")

    if "vectors" in st.session_state:
        st.caption(f"Prepared: {st.session_state.indexed_source}")

    if st.button("Clear Chat", use_container_width=True):
        st.session_state.chat_messages = []

topics_tab, ask_tab, quiz_tab, flashcard_tab = st.tabs(
    ["Topics", "Chat", "Generate Quiz", "Flashcards"]
)

with topics_tab:
    st.subheader("Topics from Table of Contents")

    if not uploaded_files:
        st.info("Upload learning material to view its topics.")
    else:
        topics_by_file, topic_errors = extract_uploaded_topics(uploaded_files)

        for filename, error in topic_errors:
            st.warning(f"Could not inspect the table of contents in {filename}: {error}")

        displayed_files = (
            topics_by_file
            if st.session_state.selected_source == "All Files"
            else {
                st.session_state.selected_source: topics_by_file.get(
                    st.session_state.selected_source,
                    [],
                )
            }
        )

        for filename, topics in displayed_files.items():
            st.markdown(f"### {filename}")
            if topics:
                st.markdown("\n".join(topic_label(topic) for topic in topics))
            else:
                st.caption("No table of contents or document headings were detected.")

with ask_tab:
    for message in st.session_state.chat_messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])

            if message["role"] == "assistant":
                if "response_time" in message:
                    st.caption(f"Response time: {message['response_time']:.2f} seconds")

                with st.expander("Source Passages"):
                    for doc in message.get("context", []):
                        st.write(source_caption(doc))
                        st.write(doc.page_content)
                        st.write("------------------------")

    user_prompt = st.chat_input("Ask about your study material")

    if user_prompt:
        st.session_state.chat_messages.append(
            {"role": "user", "content": user_prompt}
        )

        with st.chat_message("user"):
            st.markdown(user_prompt)

        if vector_database_ready():
            with st.chat_message("assistant"):
                with st.spinner("Thinking..."):
                    response = answer_chat_question(user_prompt)

                if response is not None:
                    st.markdown(response["content"])
                    st.caption(f"Response time: {response['response_time']:.2f} seconds")

                    with st.expander("Source Passages"):
                        for doc in response["context"]:
                            st.write(source_caption(doc))
                            st.write(doc.page_content)
                            st.write("------------------------")

            if response is not None:
                st.session_state.chat_messages.append(
                    {
                        "role": "assistant",
                        "content": response["content"],
                        "context": response["context"],
                        "response_time": response["response_time"],
                    }
                )

with quiz_tab:
    with st.form("quiz_form"):
        topic = st.text_input("Quiz topic", placeholder="For example: force and pressure")
        question_count = st.slider("Number of questions", 3, 10, 5)
        difficulty = st.selectbox("Difficulty", ["Beginner", "Intermediate", "Advanced"])
        quiz_submitted = st.form_submit_button("Generate Interactive Quiz")

    if quiz_submitted and topic and vector_database_ready() and ai_provider_ready():
        retriever = st.session_state.vectors.as_retriever(
            search_kwargs={"k": min(max(question_count * 2, 4), 10)}
        )

        with st.spinner("Creating your quiz..."):
            context_docs = retriever.invoke(topic)
            quiz_chain = create_stuff_documents_chain(llm, quiz_prompt)
            quiz_response = invoke_ai_chain(
                quiz_chain,
                {
                    "context": context_docs,
                    "topic": topic,
                    "difficulty": difficulty.lower(),
                    "question_count": question_count,
                },
            )
            if quiz_response is not None:
                questions = parse_quiz_response(quiz_response)
                if questions:
                    quiz_id = hashlib.sha256(
                        quiz_response.encode("utf-8")
                    ).hexdigest()[:12]
                    st.session_state.quiz_response = {
                        "content": quiz_response,
                        "questions": questions,
                        "quiz_id": quiz_id,
                        "topic": topic,
                        "context": context_docs,
                    }
                    reset_quiz_state()
                else:
                    st.error(
                        "The quiz could not be converted into interactive questions. "
                        "Please generate it again."
                    )

    if "quiz_response" in st.session_state:
        quiz_response = st.session_state.quiz_response
        st.subheader(f"Quiz: {quiz_response['topic']}")

        if quiz_response.get("questions"):
            render_interactive_quiz(quiz_response)
        else:
            st.warning("This quiz uses the previous format. Generate a new interactive quiz.")

        with st.expander("Quiz Source Passages"):
            for doc in quiz_response["context"]:
                st.write(source_caption(doc))
                st.write(doc.page_content)
                st.write("------------------------")

with flashcard_tab:
    with st.form("flashcard_form"):
        flashcard_topic = st.text_input(
            "Flashcard topic",
            placeholder="For example: crop production"
        )
        card_count = st.slider("Number of cards", 3, 12, 6)
        detail_level = st.selectbox(
            "Card style",
            ["Quick review", "Detailed study", "Exam preparation"]
        )
        flashcard_submitted = st.form_submit_button("Generate Flashcards")

    if (
        flashcard_submitted
        and flashcard_topic
        and vector_database_ready()
        and ai_provider_ready()
    ):
        retriever = st.session_state.vectors.as_retriever(
            search_kwargs={"k": min(max(card_count, 4), 10)}
        )

        with st.spinner("Preparing your flashcards..."):
            context_docs = retriever.invoke(flashcard_topic)
            flashcard_chain = create_stuff_documents_chain(llm, flashcard_prompt)
            flashcard_content = invoke_ai_chain(
                flashcard_chain,
                {
                    "context": context_docs,
                    "topic": flashcard_topic,
                    "detail_level": detail_level.lower(),
                    "card_count": card_count,
                },
            )
            if flashcard_content is not None:
                st.session_state.flashcard_response = {
                    "content": flashcard_content,
                    "topic": flashcard_topic,
                    "context": context_docs,
                }

    if "flashcard_response" in st.session_state:
        flashcard_response = st.session_state.flashcard_response
        cards = parse_flashcards(flashcard_response["content"])

        st.subheader(f"Flashcards: {flashcard_response['topic']}")

        if cards:
            deck_id = hashlib.sha256(
                flashcard_response["content"].encode("utf-8")
            ).hexdigest()[:12]
            columns = st.columns(2)
            for index, card in enumerate(cards, start=1):
                with columns[(index - 1) % 2]:
                    render_flashcard(card, index, deck_id)
        else:
            st.write(flashcard_response["content"])

        with st.expander("Flashcard Source Passages"):
            for doc in flashcard_response["context"]:
                st.write(source_caption(doc))
                st.write(doc.page_content)
                st.write("------------------------")
